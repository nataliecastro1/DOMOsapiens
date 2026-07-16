import io
import os
import re
import zipfile
from pathlib import Path
from fastapi import APIRouter, File, HTTPException, UploadFile
from fastapi.responses import FileResponse, Response

from services.uploads import UploadError, UPLOAD_DIR, list_uploads, save_upload

router = APIRouter(prefix="/api")

SUSPICIOUS_PATTERNS = [
    (r'\bdraft\b',              'Looks like a draft — the document may be incomplete.'),
    (r'\bcopy\b',               'Looks like a copy — this may not be the original file.'),
    (r'\bv\d+\b',               'Looks like a versioned file (v1, v2…) — confirm this is the final version.'),
    (r'\bversion\s*\d+\b',      'Looks like a versioned file — confirm this is the final version.'),
    (r'[\(\[]\s*\d+\s*[\)\]]',  'Filename contains a number in parentheses — may be a duplicate copy.'),
    (r'\bcopy\s*\d+\b',         'Looks like a numbered copy — this may not be the original.'),
    (r'\bwip\b',                'WIP (Work In Progress) detected — document may be incomplete.'),
    (r'\btemp\b',               'Filename contains "temp" — confirm this is the correct file.'),
    (r' - \d+$',                'Filename ends with a number — may be a duplicate copy.'),
]


def _check_name(filename: str) -> list[str]:
    name = Path(filename).stem.lower()
    warnings = []
    seen = set()
    for pattern, msg in SUSPICIOUS_PATTERNS:
        if re.search(pattern, name, re.IGNORECASE) and msg not in seen:
            warnings.append(msg)
            seen.add(msg)
    return warnings


def _first_slide_text(path: str) -> str:
    ext = Path(path).suffix.lower()
    try:
        if ext in ('.pptx', '.ppt'):
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            prs = Presentation(path)
            if not prs.slides:
                return ''
            slide = prs.slides[0]
            lines = []
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for child in shape.shapes:
                        if child.has_text_frame:
                            lines.extend(p.text.strip() for p in child.text_frame.paragraphs if p.text.strip())
                elif shape.has_text_frame:
                    lines.extend(p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip())
            return ' '.join(lines)
        elif ext == '.pdf':
            import base64, anthropic  # just read raw bytes for a quick text scan
            # Try simple text extraction without AI
            try:
                import pypdf
                with open(path, 'rb') as f:
                    reader = pypdf.PdfReader(f)
                    if reader.pages:
                        return reader.pages[0].extract_text() or ''
            except Exception:
                pass
    except Exception:
        pass
    return ''


@router.post("/uploads")
async def upload_file(file: UploadFile = File(...)):
    """Accept a source document (PPTX/PDF/XLSX) and store it on disk."""
    content = await file.read()
    try:
        return save_upload(file.filename, content)
    except UploadError as e:
        raise HTTPException(status_code=400, detail=str(e))


@router.get("/uploads")
def get_uploads():
    """List previously uploaded files, newest first."""
    return list_uploads()


@router.get("/uploads/{stored_name}/check")
def check_upload(stored_name: str, client: str = '', publisher: str = '', year: str = '', original_filename: str = ''):
    """Scan an uploaded file for red flags: draft/copy/versioned filename, and
    whether the first page text mentions the expected client, publisher, and year."""
    safe = os.path.basename(stored_name)
    path = os.path.join(UPLOAD_DIR, safe)
    if not os.path.exists(path):
        raise HTTPException(status_code=404, detail="File not found")

    # Always check the ORIGINAL filename — the stored name is a UUID hash
    name_to_check = original_filename if original_filename else safe
    warnings = _check_name(name_to_check)

    first_text = _first_slide_text(path)
    title = first_text[:120] if first_text else safe

    # Client / publisher / year presence check (all case-insensitive)
    if client and not re.search(re.escape(client), first_text, re.IGNORECASE):
        warnings.append(f'No match found for client "{client}" — this document may belong to a different client.')
    if publisher and not re.search(re.escape(publisher), first_text, re.IGNORECASE):
        warnings.append(f'No match found for publisher "{publisher}" — this document may be for a different publisher.')
    if year and not re.search(re.escape(year), first_text):
        warnings.append(f'No match found for "{year}" — the document might be from a different year.')

    return {
        'title': title,
        'warnings': warnings,
        'is_ok': len(warnings) == 0,
    }


@router.get("/uploads/{stored_name}/thumbnail")
def serve_thumbnail(stored_name: str):
    """Return the cover thumbnail for a document.
    PPTX: extracts the embedded docProps/thumbnail.jpeg (instant, zero conversion).
    PDF: renders page 1 at 2× via fitz."""
    safe = os.path.basename(stored_name)
    path = os.path.join(UPLOAD_DIR, safe)
    if not os.path.exists(path):
        raise HTTPException(status_code=404, detail="File not found")
    ext = Path(path).suffix.lower()

    if ext in ('.pptx', '.ppt'):
        with zipfile.ZipFile(path) as z:
            for candidate in ('docProps/thumbnail.jpeg', 'docProps/thumbnail.jpg', 'docProps/thumbnail.png'):
                if candidate in z.namelist():
                    mime = 'image/png' if candidate.endswith('.png') else 'image/jpeg'
                    return Response(z.read(candidate), media_type=mime)
        raise HTTPException(status_code=404, detail="No embedded thumbnail")

    if ext == '.pdf':
        try:
            import fitz
            doc = fitz.open(path)
            page = doc[0]
            pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0))
            data = pix.tobytes('png')
            doc.close()
            return Response(data, media_type='image/png')
        except Exception as e:
            raise HTTPException(status_code=500, detail=str(e))

    raise HTTPException(status_code=415, detail="Preview not supported for this file type")


@router.get("/uploads/{stored_name}/preview.pdf")
def serve_preview_pdf(stored_name: str):
    """Serve a PDF for in-browser rendering via PDF.js (PDFs only)."""
    safe = os.path.basename(stored_name)
    path = os.path.join(UPLOAD_DIR, safe)
    if not os.path.exists(path):
        raise HTTPException(status_code=404, detail="File not found")
    if Path(path).suffix.lower() != '.pdf':
        raise HTTPException(status_code=415, detail="PDF only")
    return FileResponse(path, media_type='application/pdf')


@router.get("/uploads/{stored_name}")
def serve_upload(stored_name: str):
    """Serve an uploaded file by its stored name."""
    safe = os.path.basename(stored_name)
    path = os.path.join(UPLOAD_DIR, safe)
    if not os.path.exists(path):
        raise HTTPException(status_code=404, detail="File not found")
    return FileResponse(path)
