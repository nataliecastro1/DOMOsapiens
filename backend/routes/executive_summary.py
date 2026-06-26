"""
Executive Summary endpoint — rich structured extraction from ROAR documents.
"""
import base64
import json
import os
import re
from pathlib import Path
from typing import Any, Optional

import anthropic
from fastapi import APIRouter, HTTPException
from pydantic import BaseModel

from config import ANTHROPIC_API_KEY, DOCUMENTS_DIR
from services.uploads import UPLOAD_DIR

router = APIRouter(prefix="/api")

# ─── Prompts ──────────────────────────────────────────────────────────────────

SUMMARY_PROMPT = """You are an expert IT Asset Management consultant at Anglepoint, a leading SAM advisory firm.
You are reading a ROAR (Return on Anglepoint Relationship) document.

Extract ALL available information and return ONLY valid JSON (no markdown, no explanation):
{
  "overview": "2-3 sentence engagement overview: client name, SAM scope, relationship context",
  "key_accomplishments": ["Specific accomplishment with numbers/publishers/license counts", "..."],
  "key_metrics": [
    {"label": "Metric name", "value": "$X.XM or a number", "context": "brief one-line explanation"}
  ],
  "recommendations": ["Actionable recommendation with specific details", "..."],
  "primary_risks": ["Specific identified risk with context and exposure", "..."],
  "market_risks": ["Vendor or market-specific risk", "..."],
  "additional_insights": ["Other relevant insight from the document", "..."],
  "next_steps": ["Concrete next step, with timeline if mentioned", "..."],
  "highlights": [
    {"label": "Total Identified Risk", "value": "$X.XM"},
    {"label": "Cost Avoidance Accomplished", "value": "$X.XM"}
  ],
  "charts": {
    "roi_breakdown": [{"name": "Category name", "value": 1000000}],
    "accomplishment_rate": {"accomplished": 40, "remaining": 60}
  }
}

Rules:
- Use ONLY information actually in the document — never invent data
- Empty arrays [] for sections with no data; do NOT omit any key
- charts.roi_breakdown values must be plain numbers (no $ signs)
- accomplishment_rate values are percentages that sum to 100
- Write in C-suite executive style: concise, data-driven, specific"""

AUGMENT_PROMPT = """You are an expert IT Asset Management consultant at Anglepoint.
The user has provided additional information to enrich an existing executive summary.

Your task:
1. Analyze the additional text provided
2. Classify each piece of information into the appropriate section
3. Rewrite the affected sections to integrate the new data professionally
4. Return the complete updated summary as valid JSON (same structure as the input)

Rules:
- Do NOT copy the user's raw text; rewrite it in executive style
- Convert raw numbers into metrics (e.g. "50 licencias generadas" → key_metrics entry)
- If additional info fits a chart, update charts.roi_breakdown
- Only update sections where the new info is relevant
- Return ONLY valid JSON, no markdown, no explanation"""


# ─── Models ───────────────────────────────────────────────────────────────────

class SummaryRequest(BaseModel):
    client: str = ""
    publisher: str = ""
    year: Optional[int] = None
    identified_risk: Optional[float] = None
    id_cost_avoidance: Optional[float] = None
    acc_cost_avoidance: Optional[float] = None
    id_cost_optimization: Optional[float] = None
    acc_cost_optimization: Optional[float] = None
    realized_savings: Optional[float] = None
    contract_spend: Optional[float] = None
    confidence: Optional[int] = None
    sme: Optional[str] = None
    stored_name: Optional[str] = None
    file_path: Optional[str] = None


class AugmentRequest(BaseModel):
    existing_summary: dict
    additional_text: str
    client: str = ""
    publisher: str = ""


# ─── Helpers ──────────────────────────────────────────────────────────────────

def _fmt(n):
    if not n:
        return None
    if n >= 1_000_000:
        return f"${n / 1_000_000:.1f}M"
    if n >= 1_000:
        return f"${n / 1_000:.0f}K"
    return f"${n:,.0f}"


def _extract_pptx_text(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides):
        lines = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        if lines:
            slides.append(f"[Slide {i + 1}]\n" + "\n".join(lines))
    return "\n\n---\n\n".join(slides)


def _build_content(body: SummaryRequest, file_path: str) -> list:
    ext = Path(file_path).suffix.lower()

    kpi_lines = [
        f"Client: {body.client or 'N/A'}",
        f"Publisher: {body.publisher or 'N/A'}",
        f"Year: {body.year or 'N/A'}",
        f"Identified Risk: {_fmt(body.identified_risk) or 'N/A'}",
        f"Identified Cost Avoidance: {_fmt(body.id_cost_avoidance) or 'N/A'}",
        f"Accomplished Cost Avoidance: {_fmt(body.acc_cost_avoidance) or 'N/A'}",
        f"Identified Cost Optimization: {_fmt(body.id_cost_optimization) or 'N/A'}",
        f"Accomplished Cost Optimization: {_fmt(body.acc_cost_optimization) or 'N/A'}",
        f"Realized Cost Savings: {_fmt(body.realized_savings) or 'N/A'}",
        f"Annual Contract Spend: {_fmt(body.contract_spend) or 'N/A'}",
        f"SME: {body.sme or 'N/A'}",
    ]
    kpi_block = "Extracted KPIs:\n" + "\n".join(kpi_lines)

    if ext == ".pdf":
        with open(file_path, "rb") as f:
            b64 = base64.standard_b64encode(f.read()).decode()
        return [
            {"type": "document", "source": {"type": "base64", "media_type": "application/pdf", "data": b64}},
            {"type": "text", "text": f"{kpi_block}\n\nRead the full document and generate the executive summary JSON."},
        ]
    elif ext in (".pptx", ".ppt"):
        doc_text = _extract_pptx_text(file_path)
        return [{"type": "text", "text": f"{kpi_block}\n\nDocument content:\n\n{doc_text}\n\nGenerate the executive summary JSON from the document above."}]
    else:
        return [{"type": "text", "text": f"{kpi_block}\n\nGenerate the executive summary JSON from the KPIs above."}]


def _resolve_file(stored_name: Optional[str], file_path: Optional[str]) -> Optional[str]:
    if stored_name:
        p = os.path.join(UPLOAD_DIR, os.path.basename(stored_name))
        if os.path.exists(p):
            return p
    if file_path:
        p = os.path.realpath(file_path)
        if os.path.exists(p):
            return p
    return None


def _call_claude(system: str, content: list, max_tokens: int = 2000) -> dict:
    if not ANTHROPIC_API_KEY:
        raise HTTPException(status_code=503, detail="Anthropic API key not configured.")
    ai = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)
    msg = ai.messages.create(
        model="claude-opus-4-5",
        max_tokens=max_tokens,
        system=system,
        messages=[{"role": "user", "content": content}],
    )
    raw = msg.content[0].text.strip()
    match = re.search(r"\{.*\}", raw, re.DOTALL)
    if not match:
        raise HTTPException(status_code=500, detail="Claude did not return valid JSON.")
    return json.loads(match.group())


def _ensure_structure(data: dict) -> dict:
    LIST_KEYS = ["key_accomplishments", "recommendations", "primary_risks",
                 "market_risks", "additional_insights", "next_steps", "highlights"]
    defaults: dict[str, Any] = {
        "overview": "",
        "key_metrics": [],
        "charts": {"roi_breakdown": [], "accomplishment_rate": {"accomplished": 0, "remaining": 100}},
        **{k: [] for k in LIST_KEYS},
    }
    for k, v in defaults.items():
        if k not in data:
            data[k] = v

    # Normalize key_metrics: Claude sometimes returns a dict or uses wrong key names
    if isinstance(data.get("key_metrics"), dict):
        converted = []
        for k, v in data["key_metrics"].items():
            if isinstance(v, dict):
                label   = v.get("label") or v.get("metric") or k.replace("_", " ").title()
                value   = str(v.get("value") or v.get("amount") or "")
                context = v.get("description") or v.get("context") or ""
                converted.append({"label": label, "value": value, "context": context})
            else:
                converted.append({"label": k.replace("_", " ").title(), "value": str(v), "context": ""})
        data["key_metrics"] = converted
    elif isinstance(data.get("key_metrics"), list):
        # Normalize list items: some may use "metric" key instead of "label"
        normalized = []
        for m in data["key_metrics"]:
            if isinstance(m, dict):
                label = m.get("label") or m.get("metric") or m.get("name") or "Metric"
                normalized.append({
                    "label":   label,
                    "value":   str(m.get("value") or m.get("amount") or ""),
                    "context": m.get("context") or m.get("description") or "",
                })
            else:
                normalized.append({"label": str(m), "value": "", "context": ""})
        data["key_metrics"] = normalized

    if not isinstance(data.get("charts"), dict):
        data["charts"] = defaults["charts"]
    data["charts"].setdefault("roi_breakdown", [])
    data["charts"].setdefault("accomplishment_rate", {"accomplished": 0, "remaining": 100})
    return data


# ─── Endpoints ────────────────────────────────────────────────────────────────

@router.post("/executive-summary")
async def generate_summary(body: SummaryRequest):
    file_path = _resolve_file(body.stored_name, body.file_path)

    if file_path:
        content = _build_content(body, file_path)
    else:
        kpi_lines = [
            f"Client: {body.client or 'N/A'}",
            f"Publisher: {body.publisher or 'N/A'}",
            f"Year: {body.year or 'N/A'}",
            f"Identified Risk: {_fmt(body.identified_risk) or 'N/A'}",
            f"Identified Cost Avoidance: {_fmt(body.id_cost_avoidance) or 'N/A'}",
            f"Accomplished Cost Avoidance: {_fmt(body.acc_cost_avoidance) or 'N/A'}",
            f"Identified Cost Optimization: {_fmt(body.id_cost_optimization) or 'N/A'}",
            f"Accomplished Cost Optimization: {_fmt(body.acc_cost_optimization) or 'N/A'}",
            f"Realized Cost Savings: {_fmt(body.realized_savings) or 'N/A'}",
        ]
        content = [{"type": "text", "text": "Generate an executive summary for this ROI engagement:\n\n" + "\n".join(kpi_lines)}]

    data = _call_claude(SUMMARY_PROMPT, content, max_tokens=2000)
    return _ensure_structure(data)


@router.post("/executive-summary/augment")
async def augment_summary(body: AugmentRequest):
    content = [
        {
            "type": "text",
            "text": (
                f"Client: {body.client or 'N/A'}\nPublisher: {body.publisher or 'N/A'}\n\n"
                f"Existing summary:\n{json.dumps(body.existing_summary, indent=2)}\n\n"
                f"Additional information from the user:\n{body.additional_text}\n\n"
                "Update the summary to integrate the additional information professionally. "
                "Return the complete updated summary as JSON."
            ),
        }
    ]
    data = _call_claude(AUGMENT_PROMPT, content, max_tokens=2000)
    return _ensure_structure(data)
