"""
ROAR ROI Extractor
==================
Purpose:
    Extract ROI fields from Anglepoint ROAR (.pptx) files for population
    into the DOMO ROI form or the Client ROI Tracker.

Fields Extracted:
    From Slide 1 (Cover):
        - client            : Client/account name
        - publisher         : Software publisher (e.g. IBM, VMware, Oracle)
        - month             : Presentation month, full name (e.g. "January"); None if absent
        - year              : Presentation year as int (e.g. 2026); None if absent

    From Core Properties (OPC metadata):
        - author, last_modified_by, created, modified, keywords,
          subject, title, description, category, revision, content_status,
          identifier, language, last_printed, version

    From Executive Summary section:
        - identified_risk
        - identified_cost_avoidance
        - accomplished_cost_avoidance
        - identified_cost_optimization   (may appear on a dedicated slide, not the
                                          consolidated exec summary slide)
        - accomplished_cost_optimization
        - remaining_risk
        - currency                        (inferred from value symbols, e.g. $ → USD)

Fields NOT Extracted (not available in a standard ROAR):
    - date_delivered          : Use cover date, filename, or SharePoint lastModified
    - realized_cost_savings   : Requires post-ROAR client action; not in ROAR
    - accomplished_cost_savings : Same as above
    - unit pricing / contract spend : Contained in ELP, not ROAR

Assumptions:
    - File is a valid .pptx following the Anglepoint ROAR template
    - Slide 1 contains the client name (BOX), publisher (TITLE first line), and date (BOX)
    - An "Executive Summary" section exists with at least one slide containing
      labeled ROI values in the format: LABEL: $VALUE[K/M/B]
    - Sub-capacity and full-capacity distinctions are NOT handled; a single
      value per field is assumed (flag for SME review if both are present)
    - Currency defaults to USD when no symbol is detected
    - Multi-publisher ROARs (rare) are not supported; first publisher found is used
    - File must be accessible as a local path

Dependencies:
    pip install python-pptx rapidfuzz

Usage:
    from roar_extractor import extract_roar
    result = extract_roar("path/to/ROAR.pptx")

    # Or from the command line:
    python roar_extractor.py path/to/ROAR.pptx
"""

import math
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from rapidfuzz import fuzz

# ---------------------------------------------------------------------------
# CONSTANTS
# ---------------------------------------------------------------------------

# Currency symbol → ISO 4217 code.
# Multi-char symbols MUST appear before their single-char prefix so the
# first match wins (e.g. "C$" before "$", "R$" before "R").
CURRENCY_SYMBOLS = [
    ("C$", "CAD"),
    ("A$", "AUD"),
    ("R$", "BRL"),
    ("$",  "USD"),
    ("£",  "GBP"),
    ("€",  "EUR"),
    ("R",  "ZAR"),   # South African Rand — after all $-variants
    ("¥",  "JPY"),
    ("₹",  "INR"),
]

# Value magnitude multipliers
MULTIPLIERS = {"K": 1_000, "M": 1_000_000, "B": 1_000_000_000}

# Month + year detection on ROAR cover slides.
# Month and year are detected INDEPENDENTLY so that whichever is present gets
# populated — a cover showing only "2026" yields a year with no month, and
# vice versa. Both may live in one text box ("January 2026") or in separate
# boxes; either way each is captured on its own.
_MONTH_NAMES = [
    "January", "February", "March", "April", "May", "June",
    "July", "August", "September", "October", "November", "December",
]

# Map every accepted spelling (full + 3-letter abbreviation) to its full name.
MONTH_LOOKUP = {}
for _full in _MONTH_NAMES:
    MONTH_LOOKUP[_full.lower()] = _full
    MONTH_LOOKUP[_full[:3].lower()] = _full
MONTH_LOOKUP["sept"] = "September"

# Longest spellings first so "sept" is tried before "sep", etc.
_MONTH_ALT = "|".join(sorted(MONTH_LOOKUP, key=len, reverse=True))
MONTH_RE = re.compile(rf"\b({_MONTH_ALT})\b", re.IGNORECASE)

# A standalone 4-digit calendar year (1900–2099).
YEAR_RE = re.compile(r"\b(19\d{2}|20\d{2})\b")

# ROI field label patterns — matched case-insensitively against individual
# shape text strings (NOT across shape boundaries).
FIELD_PATTERNS = {
    "identified_risk":               r"identified\s+(?:financial\s+)?risk",
    "identified_cost_avoidance":     r"identified\s+(?:cost\s+)?avoidance",
    "accomplished_cost_avoidance":   r"accomplished\s+(?:cost\s+)?avoidance",
    "identified_cost_optimization":  r"identified\s+cost\s+optimi[sz]ation",
    "accomplished_cost_optimization":r"accomplished\s+cost\s+optimi[sz]ation",
    "remaining_risk":                r"remaining\s+(?:financial\s+)?risk",
    "realized_cost_savings":         r"reali[sz]ed\s+cost\s+savings",
}

COMPILED_LABEL_RES = {
    field: re.compile(pattern, re.IGNORECASE)
    for field, pattern in FIELD_PATTERNS.items()
}

# Canonical human label per field, used for graded (fuzzy) label-fit scoring.
FIELD_LABELS = {
    "identified_risk":                "identified risk",
    "identified_cost_avoidance":      "identified cost avoidance",
    "accomplished_cost_avoidance":    "accomplished cost avoidance",
    "identified_cost_optimization":   "identified cost optimization",
    "accomplished_cost_optimization": "accomplished cost optimization",
    "remaining_risk":                 "remaining risk",
    "realized_cost_savings":          "realized cost savings",
}

# --- Confidence scoring -----------------------------------------------------
# A field value's confidence (0–100) is a weighted blend of four signals:
#   proximity   – how the value sat relative to its label (same box > adjacent)
#   source_role – role of the slide it came from (consolidated summary > other)
#   label_fit   – how cleanly the field label matched (fuzzy ratio)
#   agreement   – corroboration: share of candidates that support this value
# Weights sum to 1.0, so the score is naturally capped at 100.
W_PROXIMITY = 0.40
W_SOURCE    = 0.25
W_LABEL     = 0.20
W_AGREE     = 0.15

# Proximity by how the value was found relative to the label.
PROXIMITY_SAME_BOX = 1.0          # label and value in one text box (Phase 1)
# Adjacency (Phase 2) decays with shape distance: dist 1→0.75, 2→0.50, 3→0.25.
def _adjacency_proximity(dist: int) -> float:
    return round(1.0 - 0.25 * dist, 2)

# Slide-role weight: the consolidated exec-summary slide is the canonical source.
SOURCE_ROLE_CONSOLIDATED = 1.0
SOURCE_ROLE_OTHER        = 0.70

# Two candidate values "agree" (corroborate, rather than splitting into separate
# alternates) when they're within this relative tolerance — e.g. "$1.5M" vs
# "$1,520,000". The absolute floor keeps near-zero values from agreeing trivially.
VALUE_REL_TOL = 0.02   # 2%
VALUE_ABS_TOL = 1.0

# How many distinct candidate values to surface per field (1 primary + N-1 alts).
MAX_CANDIDATES_PER_FIELD = 3

# Matches a currency-prefixed value anywhere in a string.
# Handles: $222K  $ 222K  $1,234,567  $1.5M  £500K  €200  C$50M  R 125K
# The "R" for ZAR uses a lookahead to avoid matching "R" inside words
# like "REMAINING" or "RISK".
DOLLAR_VALUE_RE = re.compile(
    r"(C\$|A\$|R\$|\$|£|€|(?<!\w)R(?=\s*[\d,])|¥|₹)"  # currency symbol
    r"\s*([\d,]+(?:\.\d+)?)"                              # numeric value
    r"\s*([KMB]?)",                                       # optional multiplier
    re.IGNORECASE,
)

# Fuzzy match thresholds
SECTION_FUZZY_THRESHOLD = 80   # for identifying the "EXECUTIVE SUMMARY" section header
SLIDE_FUZZY_THRESHOLD   = 70   # for identifying the consolidated exec summary slide

# Patterns that mark the END of the executive summary section
NEXT_SECTION_MARKERS = [
    "recommendations", "questions", "next steps",
    "appendix", "q&a", "agenda",
]

# Design-element filter for cover slide shape skipping
_DESIGN_RE = re.compile(r"^[+\-–—|/\\*#@!\d]+$")


# ---------------------------------------------------------------------------
# HELPERS
# ---------------------------------------------------------------------------

def iter_shapes(shapes):
    """
    Recursively yield leaf shapes, descending into GROUP containers.

    PowerPoint layouts — especially exec summary cards — frequently wrap
    child shapes inside a group. A plain iteration of slide.shapes stops
    at the group container and misses all children. This generator recurses
    through every level of nesting so no text is skipped.
    """
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape


def get_slide_title(slide) -> str:
    """Return the title placeholder text for a slide, or empty string."""
    for shape in iter_shapes(slide.shapes):
        if (shape.has_text_frame
                and hasattr(shape, "is_placeholder")
                and shape.is_placeholder
                and shape.placeholder_format is not None
                and shape.placeholder_format.idx == 0):
            return shape.text_frame.text.strip()
    return ""


def get_slide_strings(slide) -> list[dict]:
    """
    Return all non-empty text entries from a slide as a flat list of dicts —
    one entry per shape text frame or table cell. Groups are recursed via
    iter_shapes. Each entry is:

        {"text": str, "shape_id": int | None, "shape_name": str | None}

    Entries are kept separate (not pre-joined) so the ROI extractor can do
    adjacency look-ahead: when a field label and its dollar value live in
    neighbouring shapes, the extractor checks the next entry rather than
    relying on a regex that spans shape boundaries. shape_id / shape_name are
    carried so an extracted value can be traced back to its shape for debugging
    (table cells inherit the table shape's id, with cell coords in the name).
    """
    entries = []
    for shape in iter_shapes(slide.shapes):
        sid = getattr(shape, "shape_id", None)
        sname = getattr(shape, "name", None)
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                entries.append({"text": t, "shape_id": sid, "shape_name": sname})
        elif shape.has_table:
            for r, row in enumerate(shape.table.rows):
                for c, cell in enumerate(row.cells):
                    t = cell.text.strip()
                    if t:
                        entries.append({
                            "text": t,
                            "shape_id": sid,
                            "shape_name": f"{sname}[r{r}c{c}]" if sname else f"[r{r}c{c}]",
                        })
    return entries


def get_slide_text_blocks(slide) -> list[dict]:
    """
    Return all non-empty text blocks as dicts with is_title flag.
    Used by cover extraction and section detection (both need is_title).
    """
    blocks = []
    for shape in iter_shapes(slide.shapes):
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        is_title = (
            hasattr(shape, "is_placeholder")
            and shape.is_placeholder
            and shape.placeholder_format is not None
            and shape.placeholder_format.idx == 0
        )
        blocks.append({"text": text, "is_title": is_title, "shape_name": shape.name})
    return blocks


def is_design_element(text: str) -> bool:
    """Return True if text is a decorative / non-content element."""
    t = text.strip()
    return (
        len(t) <= 2
        or t.startswith("©")
        or bool(_DESIGN_RE.match(t))
        or t.lower().startswith("distribution prohibited")
        or t.lower().startswith("confidential")
    )


def parse_dollar_match(match) -> tuple[float | None, str]:
    """
    Convert a DOLLAR_VALUE_RE regex match to (numeric value, ISO currency code).
    Returns (None, "USD") if the numeric part cannot be parsed.
    """
    sym_raw  = (match.group(1) or "").strip()
    num_str  = match.group(2) or ""
    mult_str = (match.group(3) or "").upper()

    currency = "USD"
    for raw_sym, iso in CURRENCY_SYMBOLS:
        if sym_raw == raw_sym or sym_raw.replace(" ", "") == raw_sym:
            currency = iso
            break

    try:
        number = float(num_str.replace(",", ""))
    except ValueError:
        return None, currency

    return number * MULTIPLIERS.get(mult_str, 1), currency


# ---------------------------------------------------------------------------
# CORE EXTRACTION: COVER SLIDE
# ---------------------------------------------------------------------------

def extract_cover_info(prs: Presentation) -> dict:
    """
    Extract client, publisher, and presentation date from Slide 1.

    Expected Slide 1 structure (Encova IBM example):
        [BOX]   Encova Insurance       ← client name
        [TITLE] IBM                    ← publisher (first line)
                Risk & Opportunity...
        [BOX]   JANUARY 2026          ← presentation date
        [BOX]   +                     ← design element, skipped
    """
    slide = prs.slides[0]
    blocks = get_slide_text_blocks(slide)

    publisher = None
    month = None   
    year = None           # int, e.g. 2026
    client_candidates = []

    for block in blocks:
        text = block["text"]

        if block["is_title"]:
            lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
            if lines:
                publisher = lines[0]
            continue

        # Detect month and year FIRST. is_design_element() treats all-digit
        # strings as decoration (its regex includes \d), so a year-only box
        # like "2026" would otherwise be discarded before it's ever read.
        month_match = MONTH_RE.search(text) if month is None else None
        year_match  = YEAR_RE.search(text) if year is None else None

        if month_match:
            month = MONTH_LOOKUP[month_match.group(1).lower()]
        if year_match:
            year = int(year_match.group(1))

        # A box that only supplied date info is not a client-name candidate.
        if month_match or year_match:
            continue

        if is_design_element(text):
            continue

        client_candidates.append(text)

    client = client_candidates[0] if client_candidates else None

    return {
        "client": client,
        "publisher": publisher,
        "month": month,
        "year": year,
    }


# ---------------------------------------------------------------------------
# CORE EXTRACTION: EXECUTIVE SUMMARY SECTION DETECTION
# ---------------------------------------------------------------------------

def find_exec_summary_slides(prs: Presentation) -> list[int]:
    """
    Return the 0-based slide indices that belong to the Executive Summary
    section, using a two-pass approach:

      Pass 1 — find the section header slide:
        A section header has very few content shapes and its title (or any
        text box) fuzzy-matches "executive summary". The number box ("02")
        that accompanies it is treated as a design element and ignored.

      Pass 2 — collect content slides until the next section header:
        Stops when a slide with sparse content whose title matches a known
        next-section marker (Recommendations, Questions, etc.) is found.

    Returns an empty list if the section cannot be located.
    """
    section_start = None

    for i, slide in enumerate(prs.slides):
        blocks = get_slide_text_blocks(slide)
        all_text = " ".join(b["text"] for b in blocks).lower()
        title = get_slide_title(slide).lower()

        # Primary: fuzzy match on the title placeholder
        if fuzz.partial_ratio("executive summary", title) >= SECTION_FUZZY_THRESHOLD:
            non_trivial = [b for b in blocks if not is_design_element(b["text"]) and len(b["text"]) > 5]
            if len(non_trivial) <= 3:
                section_start = i
                break

        # Fallback: "EXECUTIVE SUMMARY  02" style text box on a sparse slide
        if re.search(r"executive\s+summary", all_text) and len(blocks) <= 4:
            section_start = i
            break

    if section_start is None:
        return []

    section_slides = []
    for i in range(section_start + 1, len(prs.slides)):
        slide = prs.slides[i]
        blocks = get_slide_text_blocks(slide)
        non_trivial = [b for b in blocks if not is_design_element(b["text"])]

        if len(non_trivial) <= 2:
            title = get_slide_title(slide).lower()
            if any(fuzz.partial_ratio(m, title) >= 80 for m in NEXT_SECTION_MARKERS):
                break

        section_slides.append(i)

    return section_slides


# ---------------------------------------------------------------------------
# CORE EXTRACTION: ROI FIELDS
# ---------------------------------------------------------------------------

def extract_roi_fields_from_slide(slide) -> dict:
    """
    Collect ALL candidate ROI values from a single slide, per field, using a
    two-phase match. Unlike a first-match-wins scan, every label occurrence on
    the slide yields a candidate, so duplicate/conflicting values are preserved
    for downstream ranking instead of being silently dropped.

    Phase 1 — same-box match (proximity = 1.0):
        A shape whose text contains both the field label and a dollar value
        (e.g. "IDENTIFIED RISK:  $222K"). Most reliable.

    Phase 2 — adjacency look-ahead (proximity decays with distance):
        If the label box has no value, scan the next 1–3 entries for one. This
        handles split layouts where label and value sit in neighbouring shapes,
        but the value is only inferred by proximity, hence lower confidence.

    Returns:
        { field_name: [candidate, ...], ... }   (only fields with ≥1 candidate)
        where candidate = {
            "value": float, "currency": str, "raw": str,
            "proximity": float, "label_fit": float,
            "shape_id": int | None, "shape_name": str | None,
        }
    """
    entries = get_slide_strings(slide)
    candidates: dict[str, list] = {}

    for field, label_re in COMPILED_LABEL_RES.items():
        canonical = FIELD_LABELS[field]
        for i, entry in enumerate(entries):
            s = entry["text"]
            if not label_re.search(s):
                continue

            # Graded label match: how cleanly the canonical label appears here.
            label_fit = fuzz.partial_ratio(canonical, s.lower()) / 100.0

            # Phase 1: value in the same box as the label.
            val_match = DOLLAR_VALUE_RE.search(s)
            if val_match:
                value, currency = parse_dollar_match(val_match)
                if value is not None:
                    candidates.setdefault(field, []).append({
                        "value": value, "currency": currency, "raw": s.strip(),
                        "proximity": PROXIMITY_SAME_BOX, "label_fit": label_fit,
                        "shape_id": entry["shape_id"], "shape_name": entry["shape_name"],
                    })
                continue  # handled this label hit; keep scanning for more

            # Phase 2: value in an adjacent box (next 1–3 entries).
            for dist, j in enumerate(range(i + 1, min(i + 4, len(entries))), start=1):
                val_match = DOLLAR_VALUE_RE.search(entries[j]["text"])
                if val_match:
                    value, currency = parse_dollar_match(val_match)
                    if value is not None:
                        candidates.setdefault(field, []).append({
                            "value": value, "currency": currency,
                            "raw": f"{s.strip()} | {entries[j]['text'].strip()}",
                            "proximity": _adjacency_proximity(dist), "label_fit": label_fit,
                            "shape_id": entries[j]["shape_id"],
                            "shape_name": entries[j]["shape_name"],
                        })
                    break

    return candidates


def _values_agree(a: float, b: float) -> bool:
    """True if two values are within tolerance — i.e. not meaningfully different
    (e.g. "$1.5M" vs "$1,520,000") — so they corroborate rather than conflict."""
    return math.isclose(a, b, rel_tol=VALUE_REL_TOL, abs_tol=VALUE_ABS_TOL)


def _base_score(c: dict) -> float:
    """Confidence contribution from everything except cross-candidate agreement."""
    return (W_PROXIMITY * c["proximity"]
            + W_SOURCE * c["source_role"]
            + W_LABEL * c["label_fit"])


def rank_candidates(candidates: list[dict]) -> dict:
    """
    Reduce all candidates for one field to a single ranked result.

    Steps:
      1. Order candidates by base score (proximity + source role + label fit),
         strongest first. Collection order — later slides first — breaks ties,
         honouring the observation that canonical values sit toward the section
         end and front-slide repeats are often wrong.
      2. Cluster by value agreement (within tolerance), so near-equal numbers
         corroborate instead of competing. A cluster's vote count is how many
         raw matches support it.
      3. Score each cluster's representative: base score + agreement share
         (votes / total candidates), as a 0–100 confidence.
      4. Return the top cluster as the primary value, with up to
         MAX_CANDIDATES_PER_FIELD-1 distinct alternates (each retaining its
         source_slide / shape_id for debugging).
    """
    total = len(candidates)
    ordered = sorted(candidates, key=_base_score, reverse=True)  # stable

    clusters: list[dict] = []  # {"rep": candidate, "votes": int}
    for c in ordered:
        for cl in clusters:
            if cl["rep"]["currency"] == c["currency"] and _values_agree(cl["rep"]["value"], c["value"]):
                cl["votes"] += 1
                break
        else:
            clusters.append({"rep": c, "votes": 1})

    def scored(cl: dict, with_factors: bool) -> dict:
        rep, votes = cl["rep"], cl["votes"]
        agreement = votes / total
        confidence = round(100 * (_base_score(rep) + W_AGREE * agreement), 2)
        out = {
            "value":        rep["value"],
            "currency":     rep["currency"],
            "raw":          rep["raw"],
            "source_slide": rep["source_slide"],
            "shape_id":     rep["shape_id"],
            "shape_name":   rep["shape_name"],
            "confidence":   confidence,
            "votes":        votes,
        }
        if with_factors:
            out["confidence_factors"] = {
                "proximity":   round(rep["proximity"], 2),
                "source_role": rep["source_role"],
                "label_fit":   round(rep["label_fit"], 2),
                "agreement":   round(agreement, 2),
            }
        return out

    ranked = [scored(cl, with_factors=True) for cl in clusters]
    ranked.sort(key=lambda r: r["confidence"], reverse=True)  # stable

    primary = ranked[0]
    alternates = []
    for alt in ranked[1:MAX_CANDIDATES_PER_FIELD]:
        alternates.append({k: alt[k] for k in (
            "value", "currency", "raw", "source_slide", "shape_id", "shape_name", "confidence", "votes",
        )})
    primary["alternates"] = alternates
    return primary


def extract_roi_from_slides(prs: Presentation, slide_indices: list[int]) -> dict:
    """
    Extract ROI fields across the executive summary section, collecting every
    candidate value and ranking them per field by confidence.

    Rather than first-match-wins, this gathers candidates from ALL exec-summary
    slides, tags each with its slide role and source slide, then hands them to
    rank_candidates() to pick a primary value plus alternates. The consolidated
    summary slide (best title match) is weighted as the canonical source; other
    slides (e.g. a dedicated cost-optimization breakout) score slightly lower.

    Returns:
        {
          field_name: {
              "value": float, "currency": str, "raw": str,
              "source_slide": int, "shape_id": int|None, "shape_name": str|None,
              "confidence": float,            ← 0–100
              "confidence_factors": {...},    ← signal breakdown (debug)
              "votes": int,                   ← corroborating matches
              "alternates": [ {value, currency, raw, source_slide,
                               shape_id, shape_name, confidence, votes}, ... ],
          },
          ...
          "_currency": str   ← dominant currency across all fields' primaries
        }
    """
    # Identify the consolidated exec-summary slide (highest title match).
    # Scan back-to-front so a tie resolves to the later slide.
    best_idx, best_score = None, -1
    for idx in reversed(slide_indices):
        score = fuzz.partial_ratio("executive summary", get_slide_title(prs.slides[idx]).lower())
        if score > best_score:
            best_score, best_idx = score, idx
    consolidated_idx = best_idx if (best_idx is not None and best_score >= SLIDE_FUZZY_THRESHOLD) else None

    # Collect candidates from every section slide. Iterate back-to-front so that
    # later slides are encountered first and win confidence ties (see rank step 1).
    per_field: dict[str, list] = {}
    for idx in reversed(slide_indices):
        role = SOURCE_ROLE_CONSOLIDATED if idx == consolidated_idx else SOURCE_ROLE_OTHER
        for field, cands in extract_roi_fields_from_slide(prs.slides[idx]).items():
            for c in cands:
                c["source_slide"] = idx + 1
                c["source_role"] = role
            per_field.setdefault(field, []).extend(cands)

    results: dict = {}
    currency_votes: dict = {}
    for field, cands in per_field.items():
        primary = rank_candidates(cands)
        results[field] = primary
        currency_votes[primary["currency"]] = currency_votes.get(primary["currency"], 0) + 1

    results["_currency"] = max(currency_votes, key=currency_votes.get) if currency_votes else "USD"
    return results


# ---------------------------------------------------------------------------
# CORE EXTRACTION: CORE PROPERTIES
# ---------------------------------------------------------------------------

def extract_core_properties(prs: Presentation) -> dict:
    """
    Extract all 15 OPC core properties exposed by python-pptx.
    String fields are normalised to None when unset; datetime fields are
    serialised to ISO 8601 strings.
    """
    cp = prs.core_properties
    def _dt(v):
        return v.isoformat() if v else None
    return {
        "author":           cp.author or None,
        "category":         cp.category or None,
        "comments":         cp.comments or None,
        "content_status":   cp.content_status or None,
        "created":          _dt(cp.created),
        "identifier":       cp.identifier or None,
        "keywords":         cp.keywords or None,
        "language":         cp.language or None,
        "last_modified_by": cp.last_modified_by or None,
        "last_printed":     _dt(cp.last_printed),
        "modified":         _dt(cp.modified),
        "revision":         cp.revision or None,
        "subject":          cp.subject or None,
        "title":            cp.title or None,
        "version":          cp.version or None,
    }


# ---------------------------------------------------------------------------
# MAIN ENTRY POINT
# ---------------------------------------------------------------------------

def extract_roar(filepath: str) -> dict:
    """
    Full extraction pipeline for a single ROAR .pptx file.

    Returns:
    {
      "file":             str,
      "core_properties":  dict,
      "client":           str | None,
      "publisher":        str | None,
      "month":            str | None,
      "year":             int | None,
      "currency":         str,
      "roi_fields": {
          field_name: {
              "value":        float,
              "currency":     str,
              "raw":          str,   ← verbatim matched text for traceability
              "source_slide": int,   ← 1-indexed
              "shape_id":     int | None,   ← shape the value came from (debug)
              "shape_name":   str | None,
              "confidence":   float,        ← 0–100 (UI buckets into tiers)
              "confidence_factors": dict,   ← {proximity, source_role, label_fit, agreement}
              "votes":        int,          ← corroborating matches for this value
              "alternates":   [ {value, currency, raw, source_slide,
                                 shape_id, shape_name, confidence, votes}, ... ],
          }, ...
      },
      "fields_not_found": [str, ...],
      "warnings":         [str, ...],
    }
    """
    path = Path(filepath)
    warnings_out = []

    if not path.exists():
        raise FileNotFoundError(f"File not found: {filepath}")
    if path.suffix.lower() not in (".pptx", ".ppt"):
        warnings_out.append(f"Unexpected file extension: {path.suffix}")

    prs = Presentation(str(path))

    core_props = extract_core_properties(prs)

    cover = extract_cover_info(prs)
    if not cover["client"]:
        warnings_out.append("Could not extract client name from Slide 1.")
    if not cover["publisher"]:
        warnings_out.append("Could not extract publisher from Slide 1 title.")
    if not cover["month"]:
        warnings_out.append("Could not extract month from Slide 1.")
    if not cover["year"]:
        warnings_out.append("Could not extract year from Slide 1.")

    exec_indices = find_exec_summary_slides(prs)
    if not exec_indices:
        warnings_out.append("Executive Summary section not found. No ROI fields extracted.")
        roi_fields = {}
        dominant_currency = "USD"
    else:
        raw = extract_roi_from_slides(prs, exec_indices)
        dominant_currency = raw.pop("_currency", "USD")
        roi_fields = raw

    fields_not_found = [f for f in FIELD_PATTERNS if f not in roi_fields]

    return {
        "file":              path.name,
        "core_properties":   core_props,
        "client":            cover["client"],
        "publisher":         cover["publisher"],
        "month":             cover["month"],
        "year":              cover["year"],
        "currency":          dominant_currency,
        "roi_fields":        roi_fields,
        "fields_not_found":  fields_not_found,
        "warnings":          warnings_out,
    }


# ---------------------------------------------------------------------------
# CLI RUNNER
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    import json

    # if len(sys.argv) < 2:
    #     print("Usage: python roar_extractor.py <path_to_roar.pptx>")
    #     sys.exit(1)

    # print(sys.argv[1])

    # result = extract_roar(sys.argv[1])
    result = extract_roar(r'/Users/chrissi/Documents/OneMain - ROAR.pptx')
    # result = extract_roar(r"/Users/chrissi/Documents/2025 December_Encova Insurance_IBM_ROAR_v1.pptx")
    print(json.dumps(result, indent=2, default=str))
