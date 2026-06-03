"""
Read text content from local ROAR/ELP documents (PDF and PPTX).
Supports: .pdf, .pptx, .ppt
"""

import os
from pathlib import Path


def extract_text_from_pdf(file_path: str) -> str:
    """Extract all text from a PDF file using pypdf."""
    try:
        from pypdf import PdfReader
        reader = PdfReader(file_path)
        pages = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
        return "\n\n".join(pages)
    except Exception as e:
        return f"[Could not read PDF: {e}]"


def extract_text_from_pptx(file_path: str) -> str:
    """Extract all text from a PowerPoint file."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slides.append(f"[Slide {i}]\n" + "\n".join(texts))
        return "\n\n".join(slides)
    except Exception as e:
        return f"[Could not read PPTX: {e}]"


def extract_text(file_path: str) -> str:
    """Extract text from a file based on its extension."""
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext in (".pptx", ".ppt"):
        return extract_text_from_pptx(file_path)
    else:
        return f"[Unsupported file type: {ext}]"
