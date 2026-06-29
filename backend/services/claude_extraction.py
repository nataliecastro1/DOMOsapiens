"""
ROI extraction using the Anthropic Claude API directly.

TOKEN SETUP:
  Open backend/.env and make sure this line is there:
  ANTHROPIC_API_KEY=sk-ant-api03-...
"""

import base64
import json
import re
from pathlib import Path

from config import ANTHROPIC_API_KEY
from services.prompt import EXTRACTION_PROMPT


def _pdf_to_base64(file_path: str) -> str:
    with open(file_path, "rb") as f:
        return base64.standard_b64encode(f.read()).decode("utf-8")


async def extract_with_claude(file_path: str) -> dict:
    """
    Send a document file directly to Claude and return extracted ROI fields.
    Claude reads the PDF visually — works with design-heavy PowerPoints and PDFs.
    """
    if not ANTHROPIC_API_KEY:
        raise ValueError(
            "Anthropic API key is not set. "
            "Open backend/.env and add: ANTHROPIC_API_KEY=sk-ant-..."
        )

    import anthropic

    client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)

    ext = Path(file_path).suffix.lower()

    if ext == ".pdf":
        b64_data = _pdf_to_base64(file_path)
        content = [
            {
                "type": "document",
                "source": {
                    "type": "base64",
                    "media_type": "application/pdf",
                    "data": b64_data,
                },
            },
            {
                "type": "text",
                "text": "Extract the ROI data from this ROAR document and return the JSON.",
            },
        ]
    elif ext in (".pptx", ".ppt"):
        from pptx import Presentation
        from pptx.util import Pt
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        def _extract_shape_text(shape, depth=0) -> list[str]:
            """Recursively extract text from any shape, including groups and tables."""
            lines = []
            # Group shapes: recurse into children
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for child in shape.shapes:
                    lines.extend(_extract_shape_text(child, depth + 1))
                return lines
            # Tables: emit each cell with row context
            if shape.has_table:
                for row_idx, row in enumerate(shape.table.rows):
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        lines.append(" | ".join(cells))
                return lines
            # Text frames (most shapes)
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        lines.append(text)
                return lines
            # Fallback: try .text attribute
            try:
                t = shape.text.strip()
                if t:
                    lines.append(t)
            except Exception:
                pass
            return lines

        prs = Presentation(file_path)
        slides_text = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_lines = []
            for shape in slide.shapes:
                slide_lines.extend(_extract_shape_text(shape))
            if slide_lines:
                slides_text.append(f"[Slide {slide_num}]\n" + "\n".join(slide_lines))

        pptx_text = "\n\n---\n\n".join(slides_text)
        content = [
            {
                "type": "text",
                "text": f"Extract the ROI data from this ROAR PowerPoint presentation and return the JSON.\n\n{pptx_text}",
            }
        ]
    else:
        content = [
            {
                "type": "text",
                "text": "Extract the ROI data from this ROAR document and return the JSON.",
            }
        ]

    message = client.messages.create(
        model="claude-opus-4-5",
        max_tokens=1024,
        system=EXTRACTION_PROMPT,
        messages=[{"role": "user", "content": content}],
    )

    raw = message.content[0].text.strip()

    json_match = re.search(r"\{.*\}", raw, re.DOTALL)
    if not json_match:
        raise ValueError(f"Claude did not return valid JSON. Raw: {raw[:300]}")

    return json.loads(json_match.group())
